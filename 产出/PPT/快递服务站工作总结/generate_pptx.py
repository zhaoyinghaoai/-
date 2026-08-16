# -*- coding: utf-8 -*-
"""
快递服务站工作总结PPT — 标准商务风格
白底深蓝 · 标题32pt · 正文18pt · 每页≤6行
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import os

# ============ 色板 ============
DEEP_BLUE   = RGBColor(0x1F, 0x4E, 0x79)  # 主色 深蓝
MID_BLUE    = RGBColor(0x2E, 0x75, 0xB6)  # 次色 中蓝
LIGHT_BLUE  = RGBColor(0xD6, 0xE4, 0xF0)  # 浅蓝底
LIGHT_GRAY  = RGBColor(0xF2, 0xF2, 0xF2)  # 浅灰底
MID_GRAY    = RGBColor(0xBF, 0xBF, 0xBF)  # 中灰线
DARK_TEXT   = RGBColor(0x33, 0x33, 0x33)  # 正文黑
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT_RED  = RGBColor(0xC0, 0x39, 0x2B)  # 强调红
ACCENT_GREEN= RGBColor(0x27, 0x7B, 0x5D)  # 强调绿

FONT_TITLE = "微软雅黑"
FONT_BODY  = "微软雅黑"

# ============ 创建演示文稿 ============
prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

SLIDE_W = prs.slide_width
SLIDE_H = prs.slide_height
MARGIN_L = Inches(0.8)
MARGIN_R = Inches(0.8)
CONTENT_W = SLIDE_W - MARGIN_L - MARGIN_R

blank_layout = prs.slide_layouts[6]  # 空白布局

# ============ 工具函数 ============
def add_slide(bg_color=WHITE):
    slide = prs.slides.add_slide(blank_layout)
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = bg_color
    return slide

def add_rect(slide, left, top, width, height, color, line_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False
    return shape

def add_text(slide, left, top, width, height, text, font_size=18,
             font_name=FONT_BODY, color=DARK_TEXT, bold=False, alignment=PP_ALIGN.LEFT,
             line_spacing=1.5):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.name = font_name
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = alignment
    p.line_spacing = line_spacing
    return txBox

def add_multi_text(slide, left, top, width, height, lines, font_size=18,
                   font_name=FONT_BODY, color=DARK_TEXT, bold=False, alignment=PP_ALIGN.LEFT,
                   line_spacing=1.5, bullet=False):
    """多行文本，lines是列表[(text, indent_level)]"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(lines):
        if isinstance(item, tuple):
            text, level = item
        else:
            text, level = item, 0
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        prefix = "    " * level
        if bullet and level == 0:
            prefix = "· "
        p.text = prefix + text
        p.font.size = Pt(font_size)
        p.font.name = font_name
        p.font.color.rgb = color
        p.font.bold = bold
        p.alignment = alignment
        p.line_spacing = line_spacing
    return txBox

def add_title_bar(slide, title_text, subtitle_text=None):
    """标准标题栏：左侧深蓝色块 + 标题文字"""
    # 标题色块
    bar_height = Inches(0.08)
    bar = add_rect(slide, MARGIN_L, Inches(0.5), Inches(0.12), Inches(0.5), DEEP_BLUE)

    # 标题文字
    add_text(slide, Inches(1.05), Inches(0.42), Inches(10), Inches(0.7),
             title_text, font_size=32, font_name=FONT_TITLE, color=DEEP_BLUE, bold=True)

    if subtitle_text:
        add_text(slide, Inches(1.05), Inches(1.05), Inches(10), Inches(0.4),
                 subtitle_text, font_size=14, font_name=FONT_BODY, color=MID_BLUE, bold=False)

    # 底部细线
    line = add_rect(slide, MARGIN_L, Inches(1.5), CONTENT_W, Pt(1.5), MID_GRAY)

def add_footer(slide, page_num, total=10):
    """页脚"""
    add_text(slide, MARGIN_L, Inches(7.05), Inches(6), Inches(0.35),
             "快递服务站 · 2026年度工作总结", font_size=10, color=MID_GRAY)
    add_text(slide, Inches(11.5), Inches(7.05), Inches(1.5), Inches(0.35),
             f"{page_num} / {total}", font_size=10, color=MID_GRAY, alignment=PP_ALIGN.RIGHT)

def add_stat_card(slide, left, top, width, height, label, number, unit, note):
    """数据卡片：浅灰底 + 大数字"""
    # 卡片背景
    card = add_rect(slide, left, top, width, height, LIGHT_GRAY)
    # 顶部色条
    add_rect(slide, left, top, width, Inches(0.06), DEEP_BLUE)
    # 标签
    add_text(slide, left + Inches(0.2), top + Inches(0.2), width - Inches(0.4), Inches(0.3),
             label, font_size=11, color=MID_GRAY, bold=False)
    # 大数字
    add_text(slide, left + Inches(0.2), top + Inches(0.55), width - Inches(0.4), Inches(0.8),
             number, font_size=36, font_name=FONT_TITLE, color=DEEP_BLUE, bold=True)
    # 单位
    if unit:
        # 在数字后面加单位（简化处理）
        pass
    # 说明
    add_text(slide, left + Inches(0.2), top + height - Inches(0.55), width - Inches(0.4), Inches(0.4),
             note, font_size=12, color=DARK_TEXT)

def add_number_with_unit(slide, left, top, number, unit, font_size=36, color=DEEP_BLUE):
    """数字+单位（单位小一号）"""
    txBox = slide.shapes.add_textbox(left, top, Inches(3), Inches(0.8))
    tf = txBox.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    run1 = p.add_run()
    run1.text = number
    run1.font.size = Pt(font_size)
    run1.font.name = FONT_TITLE
    run1.font.color.rgb = color
    run1.font.bold = True
    run2 = p.add_run()
    run2.text = " " + unit
    run2.font.size = Pt(font_size // 2)
    run2.font.name = FONT_BODY
    run2.font.color.rgb = color
    run2.font.bold = False
    return txBox

# ================================================================
# Slide 1 · 封面
# ================================================================
slide = add_slide(DEEP_BLUE)

# 大标题
add_text(slide, Inches(1.5), Inches(2.0), Inches(10), Inches(1.2),
         "快递服务站", font_size=54, font_name=FONT_TITLE, color=WHITE, bold=True)

# 副标题
add_text(slide, Inches(1.5), Inches(3.3), Inches(10), Inches(0.8),
         "2026 年度工作总结报告", font_size=28, font_name=FONT_TITLE, color=RGBColor(0xA0, 0xC4, 0xE8), bold=False)

# 分割线
add_rect(slide, Inches(1.5), Inches(4.3), Inches(3), Pt(2), RGBColor(0xA0, 0xC4, 0xE8))

# 信息行
add_text(slide, Inches(1.5), Inches(4.6), Inches(10), Inches(0.5),
         "汇报人：赵英浩    |    汇报日期：2026年8月", font_size=16, color=RGBColor(0xCC, 0xDD, 0xEE))

# 底部装饰条
add_rect(slide, Inches(0), Inches(7.2), SLIDE_W, Inches(0.3), MID_BLUE)

# ================================================================
# Slide 2 · 核心数据
# ================================================================
slide = add_slide()
add_title_bar(slide, "核心运营数据", "全年无休 · 持续服务")

# 2x3 数据卡片
card_w = Inches(3.6)
card_h = Inches(2.2)
gap = Inches(0.3)
start_x = MARGIN_L
start_y = Inches(1.9)

stats = [
    ("日均件量", "~100", "件/天", "日均代收代发快件量"),
    ("年度总件量", "36,500", "件", "全年累计处理快件"),
    ("服务时长", "12.5", "小时/天", "早7:00 至 晚19:30"),
    ("差错记录", "0", "件", "全年零丢件零错发"),
    ("运营天数", "365", "天", "全年无休持续运营"),
    ("覆盖户数", "200+", "户", "覆盖全村常住人口"),
]

for i, (label, num, unit, note) in enumerate(stats):
    row = i // 3
    col = i % 3
    x = start_x + col * (card_w + gap)
    y = start_y + row * (card_h + gap)
    add_stat_card(slide, x, y, card_w, card_h, label, num, unit, note)

add_footer(slide, 2)

# ================================================================
# Slide 3 · 日常运营
# ================================================================
slide = add_slide()
add_title_bar(slide, "日常运营概述", "12.5小时不间断服务")

# 左侧正文
add_multi_text(slide, MARGIN_L, Inches(1.9), Inches(6.5), Inches(4.5), [
    ("每天清晨7点开门，晚上7点半收摊。快递从镇上分拨中心送到村里，", 0),
    ("登记、分拣、上架、通知取件——这套流程每天重复近100次。", 0),
    ("", 0),
    ('村级代收点的核心价值是\u201C最后一公里\u201D：', 0),
    ("村民不用赶去镇上取件", 1),
    ("快递不用因无人签收而退回", 1),
    ("每一件包裹，都在这里完成最后一步交接", 0),
], font_size=18, line_spacing=1.8, bullet=True)

# 右侧时间线
timeline_x = Inches(8.0)
timeline_w = Inches(4.5)
add_text(slide, timeline_x, Inches(1.9), timeline_w, Inches(0.4),
         "关键时间节点", font_size=16, color=DEEP_BLUE, bold=True)

times = [
    ("07:00", "开门接件", "第一批快递到达"),
    ("11-14时", "取件高峰", "村民集中来取"),
    ("19:30", "当日清件", "登记入库完毕"),
]

for i, (time, title, desc) in enumerate(times):
    y = Inches(2.5 + i * 1.5)
    # 圆点
    add_rect(slide, timeline_x, y + Inches(0.15), Inches(0.15), Inches(0.15), DEEP_BLUE)
    # 时间
    add_text(slide, timeline_x + Inches(0.35), y, Inches(1.5), Inches(0.4),
             time, font_size=20, font_name=FONT_TITLE, color=DEEP_BLUE, bold=True)
    # 标题+说明
    add_text(slide, timeline_x + Inches(0.35), y + Inches(0.45), Inches(4), Inches(0.35),
             title + " · " + desc, font_size=14, color=DARK_TEXT)
    # 分隔线
    if i < 2:
        add_rect(slide, timeline_x + Inches(0.07), y + Inches(0.4), Pt(1), Inches(1.0), MID_GRAY)

add_footer(slide, 3)

# ================================================================
# Slide 4 · 工作流程
# ================================================================
slide = add_slide()
add_title_bar(slide, "标准作业流程", "快递代收五步法")

steps = [
    ("01", "接收入库", "清点数量\n核对签收单"),
    ("02", "登记编号", "逐件录入\n生成取件码"),
    ("03", "分拣上架", "编号分区\n便于取件"),
    ("04", "通知取件", "短信/微信\n通知收件人"),
    ("05", "核验交付", "凭码取件\n签字确认"),
]

step_w = Inches(2.1)
step_h = Inches(3.0)
arrow_w = Inches(0.3)
total_w = 5 * step_w + 4 * arrow_w
start_x2 = (SLIDE_W - total_w) // 2
step_y = Inches(2.3)

for i, (num, title, desc) in enumerate(steps):
    x = start_x2 + i * (step_w + arrow_w)
    # 卡片
    add_rect(slide, x, step_y, step_w, step_h, LIGHT_GRAY)
    # 顶部色条
    add_rect(slide, x, step_y, step_w, Inches(0.08), DEEP_BLUE)
    # 编号
    add_text(slide, x, step_y + Inches(0.25), step_w, Inches(0.5),
             num, font_size=28, font_name=FONT_TITLE, color=DEEP_BLUE, bold=True, alignment=PP_ALIGN.CENTER)
    # 标题
    add_text(slide, x, step_y + Inches(0.9), step_w, Inches(0.5),
             title, font_size=18, font_name=FONT_TITLE, color=DARK_TEXT, bold=True, alignment=PP_ALIGN.CENTER)
    # 分隔线
    add_rect(slide, x + Inches(0.5), step_y + Inches(1.5), step_w - Inches(1.0), Pt(1), MID_GRAY)
    # 说明
    add_text(slide, x + Inches(0.2), step_y + Inches(1.7), step_w - Inches(0.4), Inches(1.0),
             desc, font_size=13, color=DARK_TEXT, alignment=PP_ALIGN.CENTER, line_spacing=1.5)

    # 箭头
    if i < 4:
        arrow_x = x + step_w + Inches(0.02)
        arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, arrow_x, step_y + Inches(1.2), arrow_w - Inches(0.04), Inches(0.4))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = MID_BLUE
        arrow.line.fill.background()
        arrow.shadow.inherit = False

# 底部说明
add_text(slide, MARGIN_L, Inches(5.6), CONTENT_W, Inches(0.4),
         "异常处理：滞留件（超3天未取→二次通知）/ 破损件（拍照留档→联系发件方）/ 退件（登记后交镇上取件）",
         font_size=13, color=MID_BLUE, alignment=PP_ALIGN.CENTER)

add_footer(slide, 4)

# ================================================================
# Slide 5 · 章节幕封
# ================================================================
slide = add_slide(DEEP_BLUE)

add_text(slide, Inches(1.5), Inches(2.5), Inches(10), Inches(1.0),
         "第二部分", font_size=24, color=RGBColor(0xA0, 0xC4, 0xE8), bold=False)

add_text(slide, Inches(1.5), Inches(3.2), Inches(10), Inches(1.5),
         "服务优化", font_size=60, font_name=FONT_TITLE, color=WHITE, bold=True)

add_rect(slide, Inches(1.5), Inches(4.8), Inches(3), Pt(2), RGBColor(0xA0, 0xC4, 0xE8))

add_text(slide, Inches(1.5), Inches(5.1), Inches(10), Inches(0.6),
         "数字化记录 · 主动通知 · 效率提升", font_size=20, color=RGBColor(0xCC, 0xDD, 0xEE))

# ================================================================
# Slide 6 · 服务亮点
# ================================================================
slide = add_slide()
add_title_bar(slide, "服务亮点", "从手写到数字化 · 从被动到主动")

# 3个大卡片
card_w2 = Inches(3.6)
card_h2 = Inches(3.5)
gap2 = Inches(0.4)
start_x3 = (SLIDE_W - 3 * card_w2 - 2 * gap2) // 2
card_y = Inches(2.2)

highlights = [
    ("100%", "数字化登记", "全部快件电子登记\n告别手写台账\n一键可查", ACCENT_GREEN),
    ("↓40%", "取件效率", "取件平均等候时间降低\n编号取件替代翻找\n高峰排队减少", ACCENT_RED),
    ("98%", "满意度", "村民满意度调查\n全年零投诉\n服务认可度高", ACCENT_GREEN),
]

for i, (num, title, desc, accent) in enumerate(highlights):
    x = start_x3 + i * (card_w2 + gap2)
    # 卡片
    add_rect(slide, x, card_y, card_w2, card_h2, LIGHT_GRAY)
    # 顶部色条
    add_rect(slide, x, card_y, card_w2, Inches(0.1), accent)
    # 大数字
    add_text(slide, x, card_y + Inches(0.5), card_w2, Inches(1.0),
             num, font_size=48, font_name=FONT_TITLE, color=accent, bold=True, alignment=PP_ALIGN.CENTER)
    # 标题
    add_text(slide, x, card_y + Inches(1.6), card_w2, Inches(0.5),
             title, font_size=20, font_name=FONT_TITLE, color=DEEP_BLUE, bold=True, alignment=PP_ALIGN.CENTER)
    # 分隔线
    add_rect(slide, x + Inches(1.0), card_y + Inches(2.2), card_w2 - Inches(2.0), Pt(1), MID_GRAY)
    # 说明
    add_text(slide, x + Inches(0.3), card_y + Inches(2.4), card_w2 - Inches(0.6), Inches(1.0),
             desc, font_size=14, color=DARK_TEXT, alignment=PP_ALIGN.CENTER, line_spacing=1.6)

add_footer(slide, 6)

# ================================================================
# Slide 7 · 改进前后对比
# ================================================================
slide = add_slide()
add_title_bar(slide, "改进前后对比", "从手写台账到数字化管理")

# 左列：Before
col_w = Inches(5.3)
col_h = Inches(4.5)
left_x = MARGIN_L
right_x = MARGIN_L + col_w + Inches(0.6)
col_y = Inches(1.9)

# 左列卡片
add_rect(slide, left_x, col_y, col_w, col_h, LIGHT_GRAY)
add_rect(slide, left_x, col_y, col_w, Inches(0.5), MID_GRAY)
add_text(slide, left_x, col_y + Inches(0.08), col_w, Inches(0.4),
         "Before · 传统模式", font_size=16, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

before_items = [
    "纸质本登记，字迹模糊难辨认",
    "取件靠翻找，高峰期排队等待",
    "通知靠口口相传，常有人忘取",
    "滞留件无跟踪，3天后才知道",
    "月度统计手工汇总，耗时半天",
]
for i, item in enumerate(before_items):
    y = col_y + Inches(0.8 + i * 0.65)
    add_text(slide, left_x + Inches(0.4), y, col_w - Inches(0.6), Inches(0.4),
             "✗ " + item, font_size=15, color=DARK_TEXT)

# 右列卡片
add_rect(slide, right_x, col_y, col_w, col_h, LIGHT_GRAY)
add_rect(slide, right_x, col_y, col_w, Inches(0.5), DEEP_BLUE)
add_text(slide, right_x, col_y + Inches(0.08), col_w, Inches(0.4),
         "After · 数字化模式", font_size=16, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

after_items = [
    "扫码录入系统，自动生成取件码",
    "编号分区上架，凭码秒取",
    "短信/微信自动通知到人",
    "超期自动预警，主动二次通知",
    "一键导出月报，数据实时可查",
]
for i, item in enumerate(after_items):
    y = col_y + Inches(0.8 + i * 0.65)
    add_text(slide, right_x + Inches(0.4), y, col_w - Inches(0.6), Inches(0.4),
             "✓ " + item, font_size=15, color=DARK_TEXT)

add_footer(slide, 7)

# ================================================================
# Slide 8 · 下一步挑战
# ================================================================
slide = add_slide(DEEP_BLUE)

add_text(slide, Inches(1.5), Inches(2.0), Inches(10), Inches(0.6),
         "下一步", font_size=24, color=RGBColor(0xA0, 0xC4, 0xE8))

add_text(slide, Inches(1.5), Inches(2.8), Inches(10), Inches(2.0),
         "村级快递代收，\n下一步怎么走？", font_size=44, font_name=FONT_TITLE, color=WHITE, bold=True,
         line_spacing=1.3)

add_rect(slide, Inches(1.5), Inches(5.0), Inches(3), Pt(2), RGBColor(0xA0, 0xC4, 0xE8))

add_text(slide, Inches(1.5), Inches(5.3), Inches(10), Inches(0.8),
         "快递量趋稳，但服务边界还可以更宽——\n代收之外，还能为村民做什么？",
         font_size=18, color=RGBColor(0xCC, 0xDD, 0xEE), line_spacing=1.6)

# ================================================================
# Slide 9 · 服务理念
# ================================================================
slide = add_slide()
add_title_bar(slide, "服务理念", "村级快递代收点的价值")

# 大引号
add_text(slide, Inches(2.0), Inches(2.2), Inches(2), Inches(1.5),
         '\u201C', font_size=80, font_name=FONT_TITLE, color=LIGHT_BLUE, bold=True)

# 金句
add_text(slide, Inches(3.5), Inches(2.8), Inches(7), Inches(2.0),
         "最后一步路，\n交给我来走。",
         font_size=40, font_name=FONT_TITLE, color=DEEP_BLUE, bold=True, line_spacing=1.4)

# 说明
add_text(slide, Inches(3.5), Inches(5.0), Inches(7), Inches(1.2),
         '村级快递代收点的意义，不在于量大，而在于\u201C有人在\u201D。\n全年365天，早7到晚7半，始终在这里。',
         font_size=16, color=DARK_TEXT, line_spacing=1.8)

# 出处
add_text(slide, Inches(3.5), Inches(6.2), Inches(7), Inches(0.4),
         "—— 快递服务站 · 2026年度", font_size=13, color=MID_GRAY)

add_footer(slide, 9)

# ================================================================
# Slide 10 · 致谢
# ================================================================
slide = add_slide(DEEP_BLUE)

add_text(slide, Inches(1.5), Inches(2.2), Inches(10), Inches(1.2),
         "感谢倾听", font_size=54, font_name=FONT_TITLE, color=WHITE, bold=True)

add_text(slide, Inches(1.5), Inches(3.5), Inches(10), Inches(0.8),
         "2026 年度快递服务站工作总结", font_size=24, color=RGBColor(0xA0, 0xC4, 0xE8))

add_rect(slide, Inches(1.5), Inches(4.5), Inches(3), Pt(2), RGBColor(0xA0, 0xC4, 0xE8))

add_text(slide, Inches(1.5), Inches(4.8), Inches(10), Inches(1.0),
         "全年365天 · 36,500件快递 · 日均12.5小时服务\n每一件都安全送达，每一天都坚守在这里",
         font_size=16, color=RGBColor(0xCC, 0xDD, 0xEE), line_spacing=1.8)

# 底部信息
add_text(slide, Inches(1.5), Inches(6.3), Inches(10), Inches(0.5),
         "汇报人：赵英浩    |    快递服务站    |    2026.08",
         font_size=14, color=RGBColor(0x99, 0xBB, 0xDD))

# 底部装饰条
add_rect(slide, Inches(0), Inches(7.2), SLIDE_W, Inches(0.3), MID_BLUE)

# ============ 保存 ============
output_path = r"E:\北辰\浩总的个人知识库\产出\PPT\快递服务站工作总结\快递服务站工作总结_2026年度_v1.pptx"
prs.save(output_path)
print(f"已保存：{output_path}")
print(f"共 {len(prs.slides)} 页")
